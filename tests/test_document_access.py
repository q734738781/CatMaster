from __future__ import annotations

import asyncio
from pathlib import Path
from zipfile import ZipFile

import pypdf
from docx import Document
from langchain.agents.middleware import ModelResponse
from langchain_core.messages import AIMessage, HumanMessage, ToolMessage
from openpyxl import Workbook
from pptx import Presentation
from pptx.util import Inches

from catmaster.runtime import document_access
from catmaster.runtime.checkpoint_serde import DocumentSafeCheckpointSerializer
from catmaster.runtime.document_access import (
    DocumentAccessMiddleware,
    ReadDocumentInput,
    extract_docx_text,
    extract_pdf_text,
    extract_pptx_text,
    extract_xlsx_text,
    read_document,
    sanitize_document_message,
    sanitize_document_tool_message,
)
from catmaster.webui.agent_loop import ThreadAgentLoopService


class _FakeContents:
    def __init__(self, data: bytes = b"text") -> None:
        self.data = data

    def get_data(self) -> bytes:
        return self.data


class _FakePage:
    def __init__(self, text: str) -> None:
        self.text = text

    def get_contents(self) -> _FakeContents:
        return _FakeContents()

    def extract_text(self) -> str:
        return self.text


class _FakeReader:
    def __init__(self, _path: Path, *, strict: bool = False) -> None:
        _ = strict
        self.pages = [_FakePage("alpha"), _FakePage("beta"), _FakePage("gamma")]


class _ToolRequest:
    def __init__(self, *, name: str, file_path: str, call_id: str = "call-1") -> None:
        self.tool_call = {"name": name, "id": call_id, "args": {"file_path": file_path}}


class _ModelRequest:
    def __init__(self, messages: list[object]) -> None:
        self.messages = messages

    def override(self, **kwargs):
        return _ModelRequest(list(kwargs.get("messages", self.messages)))


def _pdf_tool_message() -> ToolMessage:
    return ToolMessage(
        content_blocks=[
            {
                "type": "file",
                "base64": "JVBERi0xLjQ=",
                "mime_type": "application/pdf",
            }
        ],
        additional_kwargs={
            "read_file_path": "/literature/paper.pdf",
            "read_file_media_type": "application/pdf",
        },
        tool_call_id="call-1",
        name="read_file",
        status="success",
    )


def _docx_tool_message() -> ToolMessage:
    mime = "application/vnd.openxmlformats-officedocument.wordprocessingml.document"
    return ToolMessage(
        content_blocks=[
            {
                "type": "file",
                "base64": "UEsDBAoAAAA=",
                "mime_type": mime,
                "filename": "report.docx",
            }
        ],
        additional_kwargs={
            "read_file_path": "/reports/report.docx",
            "read_file_media_type": mime,
        },
        tool_call_id="call-docx",
        name="read_file",
        status="success",
    )


def test_read_document_schema_uses_non_nullable_optional_controls(tmp_path: Path) -> None:
    schema = ReadDocumentInput.model_json_schema()
    tool_schema = DocumentAccessMiddleware(files_root=tmp_path).tools[0].args

    pages = schema["properties"]["pages"]
    assert pages["type"] == "string"
    assert pages["default"] == ""
    assert "anyOf" not in pages
    assert "pages" not in schema["required"]
    assert tool_schema["pages"]["type"] == "string"
    assert tool_schema["pages"]["default"] == ""
    assert "anyOf" not in tool_schema["pages"]
    cursor = schema["properties"]["cursor"]
    assert cursor["type"] == "string"
    assert cursor["default"] == ""
    assert "cursor" not in schema["required"]
    assert tool_schema["cursor"]["type"] == "string"
    assert tool_schema["cursor"]["default"] == ""


def test_extract_pdf_text_is_bounded_by_requested_pages(tmp_path: Path, monkeypatch) -> None:
    source = tmp_path / "paper.pdf"
    source.write_bytes(b"%PDF-fake")
    monkeypatch.setattr(pypdf, "PdfReader", _FakeReader)

    result = extract_pdf_text(source, pages="2-3")

    assert result.total_units == 3
    assert result.start_unit == 1
    assert result.end_unit == 2
    assert "--- Page 2 ---\nbeta" in result.text
    assert "--- Page 3 ---\ngamma" in result.text
    assert "alpha" not in result.text


def test_read_document_returns_pdf_content_without_recursive_read_instruction(tmp_path: Path, monkeypatch) -> None:
    source = tmp_path / "paper.pdf"
    source.write_bytes(b"%PDF-fake")
    monkeypatch.setattr(pypdf, "PdfReader", _FakeReader)

    result = read_document(tmp_path, file_path="/paper.pdf", pages="1-2")

    assert "PDF source: `/paper.pdf`" in result
    assert "alpha" in result
    assert "beta" in result
    assert "base64" not in result
    assert "read_file" not in result


def test_docx_parser_preserves_paragraphs_and_tables(tmp_path: Path) -> None:
    source = tmp_path / "report.docx"
    document = Document()
    document.add_paragraph("Catalyst summary")
    table = document.add_table(rows=2, cols=2)
    table.cell(0, 0).text = "sample"
    table.cell(0, 1).text = "yield"
    table.cell(1, 0).text = "Pd-TWC"
    table.cell(1, 1).text = "92%"
    document.save(source)

    extraction = extract_docx_text(source)
    result = read_document(tmp_path, file_path="/report.docx")

    assert extraction.truncated is False
    assert "Catalyst summary" in extraction.text
    assert "sample\tyield" in extraction.text
    assert "Pd-TWC\t92%" in result
    assert "read_file" not in result


def test_xlsx_parser_reads_workbook_in_read_only_value_mode(tmp_path: Path) -> None:
    source = tmp_path / "results.xlsx"
    workbook = Workbook()
    sheet = workbook.active
    sheet.title = "Kinetics"
    sheet.append(["temperature", "rate"])
    sheet.append([650, 1.25])
    workbook.save(source)
    workbook.close()

    extraction = extract_xlsx_text(source)
    result = read_document(tmp_path, file_path="/results.xlsx")

    assert extraction.truncated is False
    assert "--- Sheet: Kinetics ---" in extraction.text
    assert "temperature\trate" in extraction.text
    assert "650\t1.25" in result


def test_docx_and_xlsx_cursors_make_every_unit_reachable_without_loss(
    tmp_path: Path,
) -> None:
    docx_path = tmp_path / "long.docx"
    document = Document()
    for index in range(18):
        document.add_paragraph(
            f"Paragraph {index:02d}: " + (f"detail-{index} " * 12)
        )
    document.save(docx_path)

    xlsx_path = tmp_path / "long.xlsx"
    workbook = Workbook()
    sheet = workbook.active
    sheet.title = "LongData"
    for index in range(30):
        sheet.append([index, f"observation-{index}-" + ("value " * 10)])
    workbook.save(xlsx_path)
    workbook.close()

    for extractor, source, virtual_path in (
        (extract_docx_text, docx_path, "/long.docx"),
        (extract_xlsx_text, xlsx_path, "/long.xlsx"),
    ):
        complete = extractor(
            source,
            virtual_path=virtual_path,
            max_chars=1_000_000,
        )
        cursor = ""
        chunks: list[str] = []
        seen_cursors: set[str] = set()
        paths: list[str] = []
        while True:
            page = extractor(
                source,
                virtual_path=virtual_path,
                cursor=cursor,
                max_chars=97,
            )
            chunks.append(page.text)
            paths.extend([page.start_unit_path, page.end_unit_path])
            if not page.next_cursor:
                break
            assert page.next_cursor not in seen_cursors
            seen_cursors.add(page.next_cursor)
            cursor = page.next_cursor
        assert "".join(chunks) == complete.text
        assert any("paragraph:" in path or "sheet:" in path for path in paths)


def test_pptx_parser_reads_slide_text_and_tables_with_slide_ranges(tmp_path: Path) -> None:
    source = tmp_path / "review.pptx"
    presentation = Presentation()
    first = presentation.slides.add_slide(presentation.slide_layouts[6])
    first.shapes.add_textbox(Inches(1), Inches(1), Inches(4), Inches(1)).text = "Overview"
    second = presentation.slides.add_slide(presentation.slide_layouts[6])
    table = second.shapes.add_table(2, 2, Inches(1), Inches(1), Inches(5), Inches(2)).table
    table.cell(0, 0).text = "system"
    table.cell(0, 1).text = "energy"
    table.cell(1, 0).text = "O2"
    table.cell(1, 1).text = "-1.2 eV"
    presentation.save(source)

    extraction = extract_pptx_text(source, pages="2")
    result = read_document(tmp_path, file_path="/review.pptx", pages="2")

    assert extraction.start_unit == 1
    assert extraction.end_unit == 1
    assert "Overview" not in extraction.text
    assert "system\tenergy" in extraction.text
    assert "O2\t-1.2 eV" in result


def test_pages_is_rejected_for_docx_and_xlsx(tmp_path: Path) -> None:
    source = tmp_path / "report.docx"
    document = Document()
    document.add_paragraph("text")
    document.save(source)

    result = read_document(tmp_path, file_path="/report.docx", pages="2")

    assert result.startswith("Error reading document:")
    assert "pages applies only" in result


def test_office_archive_preflight_rejects_excessive_members(tmp_path: Path, monkeypatch) -> None:
    source = tmp_path / "oversized.docx"
    with ZipFile(source, "w") as archive:
        archive.writestr("one.xml", "one")
        archive.writestr("two.xml", "two")
    monkeypatch.setattr(document_access, "MAX_OFFICE_ARCHIVE_MEMBERS", 1)

    result = read_document(tmp_path, file_path="/oversized.docx")

    assert result.startswith("Error reading document:")
    assert "too many internal files" in result


def test_webui_pdf_attachment_is_sent_as_bounded_text(tmp_path: Path, monkeypatch) -> None:
    source = tmp_path / "files" / "attachments" / "paper.pdf"
    source.parent.mkdir(parents=True)
    source.write_bytes(b"%PDF-fake")
    monkeypatch.setattr(pypdf, "PdfReader", _FakeReader)
    service = object.__new__(ThreadAgentLoopService)
    service.workspace = tmp_path
    warnings: list[str] = []

    block = service._document_attachment_block(
        source,
        filename="paper.pdf",
        workspace_path="files/attachments/paper.pdf",
        warnings=warnings,
    )

    assert block is not None
    assert block["type"] == "text"
    assert "alpha" in block["text"]
    assert "base64" not in block["text"]
    assert warnings == []


def test_document_middleware_blocks_pdf_read_file_before_handler(tmp_path: Path) -> None:
    middleware = DocumentAccessMiddleware(files_root=tmp_path)
    handler_called = False

    async def _handler(_request):
        nonlocal handler_called
        handler_called = True
        return _pdf_tool_message()

    result = asyncio.run(
        middleware.awrap_tool_call(
            _ToolRequest(name="read_file", file_path="/literature/paper.pdf"),
            _handler,
        )
    )

    assert handler_called is False
    assert isinstance(result, ToolMessage)
    assert result.status == "error"
    assert "read_document" in str(result.content)
    assert "Do not retry `read_file`" in str(result.content)
    assert "base64" not in str(result.content)


def test_document_middleware_blocks_office_read_file_before_handler(tmp_path: Path) -> None:
    middleware = DocumentAccessMiddleware(files_root=tmp_path)
    handler_called = False

    async def _handler(_request):
        nonlocal handler_called
        handler_called = True
        return _docx_tool_message()

    result = asyncio.run(
        middleware.awrap_tool_call(
            _ToolRequest(name="read_file", file_path="/reports/report.docx", call_id="call-docx"),
            _handler,
        )
    )

    assert handler_called is False
    assert isinstance(result, ToolMessage)
    assert result.status == "error"
    assert "read_document" in str(result.content)


def test_office_tool_message_sanitizer_removes_replayed_binary() -> None:
    original = _docx_tool_message()

    sanitized = sanitize_document_tool_message(original)

    assert sanitized.tool_call_id == original.tool_call_id
    assert "UEsDB" not in str(sanitized.content)
    assert sanitized.additional_kwargs["read_file_path"] == "/reports/report.docx"
    assert sanitized.additional_kwargs["catmaster_document_payload_removed"] is True


def test_document_middleware_preserves_image_read_file_result(tmp_path: Path) -> None:
    middleware = DocumentAccessMiddleware(files_root=tmp_path)
    image = ToolMessage(
        content_blocks=[{"type": "image", "base64": "aW1hZ2U=", "mime_type": "image/png"}],
        additional_kwargs={"read_file_path": "/figure.png", "read_file_media_type": "image/png"},
        tool_call_id="call-image",
        name="read_file",
    )

    async def _handler(_request):
        return image

    result = asyncio.run(
        middleware.awrap_tool_call(
            _ToolRequest(name="read_file", file_path="/figure.png", call_id="call-image"),
            _handler,
        )
    )

    assert result is image
    assert result.content[0]["base64"] == "aW1hZ2U="


def test_document_tool_message_sanitizer_preserves_message_identity() -> None:
    original = _pdf_tool_message()

    sanitized = sanitize_document_tool_message(original)

    assert sanitized.id == original.id
    assert sanitized.tool_call_id == original.tool_call_id
    assert sanitized.name == "read_file"
    assert "base64" not in str(sanitized.content)
    assert sanitized.additional_kwargs["read_file_path"] == "/literature/paper.pdf"
    assert sanitized.additional_kwargs["catmaster_document_payload_removed"] is True


def test_checkpoint_serializer_removes_document_bytes_before_persistence() -> None:
    original = _pdf_tool_message()
    serializer = DocumentSafeCheckpointSerializer()

    type_name, payload = serializer.dumps_typed(
        {"channel_values": {"messages": [original]}}
    )
    restored = serializer.loads_typed((type_name, payload))
    message = restored["channel_values"]["messages"][0]

    assert b"JVBERi0xLjQ" not in payload
    assert isinstance(message, ToolMessage)
    assert message.tool_call_id == original.tool_call_id
    assert message.additional_kwargs["catmaster_document_payload_removed"] is True
    assert "read_document" in str(message.content)


def test_checkpoint_serializer_preserves_inline_images() -> None:
    original = ToolMessage(
        content_blocks=[
            {"type": "image", "base64": "aW1hZ2U=", "mime_type": "image/png"}
        ],
        additional_kwargs={
            "read_file_path": "/figure.png",
            "read_file_media_type": "image/png",
        },
        tool_call_id="call-image",
        name="read_file",
    )
    serializer = DocumentSafeCheckpointSerializer()

    restored = serializer.loads_typed(serializer.dumps_typed(original))

    assert isinstance(restored, ToolMessage)
    assert restored.content[0]["base64"] == "aW1hZ2U="


def test_document_message_sanitizer_preserves_text_and_image_blocks() -> None:
    original = HumanMessage(
        content=[
            {"type": "text", "text": "Inspect the attached evidence."},
            {"type": "file", "base64": "JVBERi0=", "mime_type": "application/pdf", "filename": "paper.pdf"},
            {"type": "image", "base64": "aW1hZ2U=", "mime_type": "image/png"},
        ]
    )

    sanitized = sanitize_document_message(original)

    assert sanitized.content[0] == original.content[0]
    assert sanitized.content[1]["type"] == "text"
    assert "paper.pdf" in sanitized.content[1]["text"]
    assert sanitized.content[2] == original.content[2]
    assert "JVBERi0=" not in str(sanitized.content)


def test_model_call_sanitizes_replayed_pdf_before_provider(tmp_path: Path) -> None:
    middleware = DocumentAccessMiddleware(files_root=tmp_path)
    request = _ModelRequest([_pdf_tool_message(), AIMessage(content="continue")])
    seen: list[object] = []

    async def _handler(sanitized_request):
        seen.extend(sanitized_request.messages)
        return ModelResponse(result=[AIMessage(content="ok")])

    response = asyncio.run(middleware.awrap_model_call(request, _handler))

    assert isinstance(response, ModelResponse)
    assert isinstance(seen[0], ToolMessage)
    assert "base64" not in str(seen[0].content)
    assert seen[1].content == "continue"
