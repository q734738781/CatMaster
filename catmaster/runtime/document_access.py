from __future__ import annotations

import asyncio
import base64
import binascii
import hashlib
import json
import mimetypes
from dataclasses import dataclass
from pathlib import Path, PurePosixPath
from typing import Any, Iterator
from zipfile import BadZipFile, ZipFile

from langchain.agents.middleware import AgentMiddleware
from langchain_core.messages import BaseMessage, ToolMessage
from langchain_core.tools import StructuredTool
from pydantic import BaseModel, ConfigDict, Field

SUPPORTED_DOCUMENT_SUFFIXES = frozenset({".pdf", ".docx", ".xlsx", ".pptx"})
DOCUMENT_MIME_TYPES = frozenset(
    {
        "application/pdf",
        "application/vnd.openxmlformats-officedocument.wordprocessingml.document",
        "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
        "application/vnd.openxmlformats-officedocument.presentationml.presentation",
    }
)

MAX_DOCUMENT_FILE_SIZE_BYTES = 50 * 1024 * 1024
MAX_DOCUMENT_TEXT_CHARS = 60_000
MAX_DOCUMENT_UNITS_PER_READ = 20
MAX_PDF_PAGE_STREAM_BYTES = 32 * 1024 * 1024
MAX_OFFICE_ARCHIVE_MEMBERS = 10_000
MAX_OFFICE_UNCOMPRESSED_BYTES = 256 * 1024 * 1024
MAX_OFFICE_MEMBER_BYTES = 128 * 1024 * 1024
MAX_XLSX_COLUMNS_PER_READ = 256


class ReadDocumentInput(BaseModel):
    """Read bounded text from a workspace document without adding binary bytes to history."""

    model_config = ConfigDict(extra="forbid")

    file_path: str = Field(
        description=(
            "Absolute virtual workspace path to a PDF, DOCX, XLSX, or PPTX file, "
            "for example /literature/paper.pdf."
        ),
    )
    pages: str = Field(
        default="",
        description=(
            "For PDF or PPTX only, an optional 1-based page/slide or range such as 3 or 3-8. "
            "Leave empty to start at the beginning; at most 20 pages or slides are returned. "
            "Leave empty for DOCX and XLSX."
        ),
    )
    cursor: str = Field(
        default="",
        description=(
            "For DOCX or XLSX only, the opaque continuation cursor returned by "
            "the preceding read. Leave empty for the first unit and for PDF/PPTX."
        ),
    )


@dataclass(frozen=True, slots=True)
class BoundedTextExtraction:
    text: str
    truncated: bool


@dataclass(frozen=True, slots=True)
class PagedTextExtraction(BoundedTextExtraction):
    total_units: int
    start_unit: int
    end_unit: int


@dataclass(frozen=True, slots=True)
class CursorTextExtraction(BoundedTextExtraction):
    next_cursor: str
    start_unit_path: str
    end_unit_path: str


class DocumentReadError(ValueError):
    """Raised when a document cannot be read within the parser safety bounds."""


class _TextCollector:
    def __init__(self, limit: int) -> None:
        self.limit = max(0, int(limit))
        self.parts: list[str] = []
        self.length = 0
        self.truncated = False

    def add(self, text: str, *, separator: str = "") -> bool:
        if not text:
            return True
        chunk = (separator if self.parts else "") + text
        remaining = self.limit - self.length
        if len(chunk) > remaining:
            if remaining > 0:
                self.parts.append(chunk[:remaining])
                self.length += remaining
            self.truncated = True
            return False
        self.parts.append(chunk)
        self.length += len(chunk)
        return True

    def render(self) -> str:
        return "".join(self.parts)


def _inspect_document_file(path: Path) -> Path:
    source = Path(path).expanduser().resolve()
    if not source.is_file():
        raise DocumentReadError(f"Document file not found: {source}")
    try:
        size_bytes = source.stat().st_size
    except OSError as exc:
        raise DocumentReadError(f"Could not inspect document file: {exc}") from exc
    if size_bytes > MAX_DOCUMENT_FILE_SIZE_BYTES:
        raise DocumentReadError(
            f"Document exceeds the {MAX_DOCUMENT_FILE_SIZE_BYTES // (1024 * 1024)} MB reading limit."
        )
    return source


def _validate_office_archive(path: Path) -> None:
    """Reject unsafe OOXML containers before an XML parser opens them."""
    try:
        with ZipFile(path) as archive:
            members = archive.infolist()
    except (BadZipFile, OSError) as exc:
        raise DocumentReadError(f"Invalid Office document: {exc}") from exc

    if len(members) > MAX_OFFICE_ARCHIVE_MEMBERS:
        raise DocumentReadError(f"Office document contains too many internal files ({len(members)}).")

    total_size = 0
    for member in members:
        if member.flag_bits & 0x1:
            raise DocumentReadError("Encrypted Office documents are not supported.")
        if member.file_size > MAX_OFFICE_MEMBER_BYTES:
            raise DocumentReadError("Office document contains an oversized internal file.")
        total_size += member.file_size
        if total_size > MAX_OFFICE_UNCOMPRESSED_BYTES:
            limit_mb = MAX_OFFICE_UNCOMPRESSED_BYTES // (1024 * 1024)
            raise DocumentReadError(f"Office document expands beyond the {limit_mb} MB safety limit.")


def _parse_unit_range(pages: str, total_units: int, *, unit_name: str) -> tuple[int, int]:
    raw = str(pages or "").strip()
    if not raw:
        return 0, total_units - 1
    values = raw.split("-")
    if len(values) not in {1, 2}:
        raise DocumentReadError(f"Invalid {unit_name} range {raw!r}; use a number or range such as 3 or 3-8.")
    try:
        start = int(values[0])
        end = int(values[-1])
    except ValueError as exc:
        raise DocumentReadError(
            f"Invalid {unit_name} range {raw!r}; use a number or range such as 3 or 3-8."
        ) from exc
    if start < 1 or end < start or start > total_units:
        raise DocumentReadError(f"{unit_name.title()} range {raw!r} is outside the document's {total_units} {unit_name}s.")
    return start - 1, min(end, total_units) - 1


def extract_pdf_text(
    path: Path,
    *,
    pages: str = "",
    max_pages: int = MAX_DOCUMENT_UNITS_PER_READ,
    max_chars: int = MAX_DOCUMENT_TEXT_CHARS,
) -> PagedTextExtraction:
    """Extract a bounded PDF page range without retaining the full document text."""
    source = _inspect_document_file(path)
    try:
        from pypdf import PdfReader

        reader = PdfReader(source, strict=False)
    except Exception as exc:
        raise DocumentReadError(f"Could not open PDF: {exc}") from exc

    total_pages = len(reader.pages)
    if total_pages == 0:
        return PagedTextExtraction("", False, 0, 0, -1)
    start, requested_end = _parse_unit_range(pages, total_pages, unit_name="page")
    end = min(requested_end, start + max(1, int(max_pages)) - 1)

    collector = _TextCollector(max_chars)
    actual_end = start - 1
    for index in range(start, end + 1):
        page = reader.pages[index]
        try:
            contents = page.get_contents()
            if contents is not None and len(contents.get_data()) > MAX_PDF_PAGE_STREAM_BYTES:
                raise DocumentReadError(
                    f"PDF page {index + 1} content stream exceeds the "
                    f"{MAX_PDF_PAGE_STREAM_BYTES // (1024 * 1024)} MB per-page limit."
                )
            page_text = (page.extract_text() or "").strip()
        except DocumentReadError:
            raise
        except Exception as exc:
            raise DocumentReadError(f"Could not extract PDF page {index + 1}: {exc}") from exc

        actual_end = index
        if page_text and not collector.add(f"--- Page {index + 1} ---\n{page_text}", separator="\n\n"):
            break

    return PagedTextExtraction(
        collector.render(),
        collector.truncated or actual_end < requested_end,
        total_pages,
        start,
        actual_end,
    )


def _docx_table_text(table: Any) -> Iterator[tuple[int, str]]:
    for row_index, row in enumerate(table.rows, start=1):
        cells = [" ".join(str(cell.text or "").splitlines()).strip() for cell in row.cells]
        line = "\t".join(cells).rstrip()
        if line.strip():
            yield row_index, line


def _document_cursor(
    *,
    virtual_path: str,
    kind: str,
    state: dict[str, int],
) -> str:
    payload = {
        "path_hash": hashlib.sha256(virtual_path.encode("utf-8")).hexdigest(),
        "kind": kind,
        "state": {key: int(value) for key, value in state.items()},
    }
    raw = json.dumps(payload, separators=(",", ":")).encode("utf-8")
    return base64.urlsafe_b64encode(raw).decode("ascii").rstrip("=")


def _document_cursor_state(
    *,
    cursor: str,
    virtual_path: str,
    kind: str,
    defaults: dict[str, int],
) -> dict[str, int]:
    value = str(cursor or "").strip()
    if not value:
        return dict(defaults)
    try:
        padded = value + "=" * (-len(value) % 4)
        payload = json.loads(base64.urlsafe_b64decode(padded).decode("utf-8"))
        if payload.get("kind") != kind:
            raise ValueError
        if payload.get("path_hash") != hashlib.sha256(
            virtual_path.encode("utf-8")
        ).hexdigest():
            raise ValueError
        raw_state = payload.get("state")
        if not isinstance(raw_state, dict):
            raise ValueError
        state = {key: int(raw_state[key]) for key in defaults}
        if any(value < 0 for value in state.values()):
            raise ValueError
        return state
    except (
        binascii.Error,
        KeyError,
        TypeError,
        UnicodeDecodeError,
        ValueError,
        json.JSONDecodeError,
    ) as exc:
        raise DocumentReadError(
            "Document continuation cursor is invalid for this file and format."
        ) from exc


def _collect_cursor_units(
    units: Iterator[tuple[dict[str, int], str, str]],
    *,
    virtual_path: str,
    kind: str,
    start_char: int,
    max_chars: int,
) -> CursorTextExtraction:
    collector = _TextCollector(max(1, int(max_chars)))
    first_path = ""
    last_path = ""
    next_cursor = ""
    first = True
    for state, unit_path, unit_text in units:
        char_offset = start_char if first else 0
        first = False
        if char_offset > len(unit_text):
            raise DocumentReadError("Document continuation cursor is outside its unit.")
        remaining_text = unit_text[char_offset:]
        if not first_path:
            first_path = (
                f"{unit_path}@{char_offset}" if char_offset else unit_path
            )
        last_path = unit_path
        available = max(0, collector.limit - collector.length)
        if len(remaining_text) > available:
            if available:
                collector.add(remaining_text[:available])
            collector.truncated = True
            next_cursor = _document_cursor(
                virtual_path=virtual_path,
                kind=kind,
                state={**state, "char": char_offset + available},
            )
            break
        collector.add(remaining_text)
    return CursorTextExtraction(
        # Do not strip page-boundary whitespace: the opaque cursor advances by
        # exact character count, so trimming here would make that content
        # permanently unreachable when pages are concatenated.
        text=collector.render(),
        truncated=bool(next_cursor),
        next_cursor=next_cursor,
        start_unit_path=first_path,
        end_unit_path=last_path,
    )


def extract_docx_text(
    path: Path,
    *,
    cursor: str = "",
    virtual_path: str = "",
    max_chars: int = MAX_DOCUMENT_TEXT_CHARS,
) -> CursorTextExtraction:
    """Extract a continuable page of DOCX paragraphs and table rows."""
    source = _inspect_document_file(path)
    cursor_path = virtual_path or str(source)
    _validate_office_archive(source)
    try:
        from docx import Document
        from docx.table import Table

        document = Document(source)
    except Exception as exc:
        raise DocumentReadError(f"Could not open DOCX: {exc}") from exc

    state = _document_cursor_state(
        cursor=cursor,
        virtual_path=cursor_path,
        kind="docx",
        defaults={"unit": 0, "char": 0},
    )

    def units() -> Iterator[tuple[dict[str, int], str, str]]:
        unit_index = 0
        paragraph_index = 0
        table_index = 0
        for block in document.iter_inner_content():
            if isinstance(block, Table):
                table_index += 1
                for row_index, line in _docx_table_text(block):
                    if unit_index >= state["unit"]:
                        yield (
                            {"unit": unit_index},
                            f"table:{table_index}/row:{row_index}",
                            f"--- Table {table_index}, row {row_index} ---\n{line}\n\n",
                        )
                    unit_index += 1
                continue
            paragraph_index += 1
            text = str(getattr(block, "text", "") or "").strip()
            if not text:
                continue
            if unit_index >= state["unit"]:
                yield (
                    {"unit": unit_index},
                    f"paragraph:{paragraph_index}",
                    f"--- Paragraph {paragraph_index} ---\n{text}\n\n",
                )
            unit_index += 1

    return _collect_cursor_units(
        units(),
        virtual_path=cursor_path,
        kind="docx",
        start_char=state["char"],
        max_chars=max_chars,
    )


def extract_xlsx_text(
    path: Path,
    *,
    cursor: str = "",
    virtual_path: str = "",
    max_chars: int = MAX_DOCUMENT_TEXT_CHARS,
) -> CursorTextExtraction:
    """Extract a continuable page of XLSX row/column units."""
    source = _inspect_document_file(path)
    cursor_path = virtual_path or str(source)
    _validate_office_archive(source)
    try:
        from openpyxl import load_workbook

        workbook = load_workbook(source, read_only=True, data_only=True)
    except Exception as exc:
        raise DocumentReadError(f"Could not open XLSX: {exc}") from exc

    state = _document_cursor_state(
        cursor=cursor,
        virtual_path=cursor_path,
        kind="xlsx",
        defaults={"unit": 0, "char": 0},
    )

    def units() -> Iterator[tuple[dict[str, int], str, str]]:
        unit_index = 0
        for sheet_name in workbook.sheetnames:
            worksheet = workbook[sheet_name]
            max_column = int(worksheet.max_column or 1)
            for row_index, row in enumerate(
                worksheet.iter_rows(
                    min_row=1,
                    max_row=int(worksheet.max_row or 1),
                    min_col=1,
                    max_col=max_column,
                    values_only=True,
                ),
                start=1,
            ):
                for start_column in range(
                    1,
                    max_column + 1,
                    MAX_XLSX_COLUMNS_PER_READ,
                ):
                    end_column = min(
                        max_column,
                        start_column + MAX_XLSX_COLUMNS_PER_READ - 1,
                    )
                    values = row[start_column - 1 : end_column]
                    row_text = "\t".join(
                        str(cell) if cell is not None else "" for cell in values
                    ).rstrip()
                    if not row_text.strip():
                        continue
                    if unit_index >= state["unit"]:
                        yield (
                            {"unit": unit_index},
                            (
                                f"sheet:{sheet_name}/row:{row_index}/"
                                f"columns:{start_column}-{end_column}"
                            ),
                            (
                                f"--- Sheet: {sheet_name} ---\n"
                                f"Unit: sheet:{sheet_name}/row:{row_index}/"
                                f"columns:{start_column}-{end_column}\n"
                                f"{row_text}\n\n"
                            ),
                        )
                    unit_index += 1

    try:
        return _collect_cursor_units(
            units(),
            virtual_path=cursor_path,
            kind="xlsx",
            start_char=state["char"],
            max_chars=max_chars,
        )
    finally:
        workbook.close()


def _iter_pptx_shape_text(shape: Any) -> Iterator[str]:
    sub_shapes = getattr(shape, "shapes", None)
    if sub_shapes is not None:
        for sub_shape in sub_shapes:
            yield from _iter_pptx_shape_text(sub_shape)
        return
    if getattr(shape, "has_table", False):
        for row in shape.table.rows:
            cells = [" ".join(str(cell.text or "").splitlines()).strip() for cell in row.cells]
            line = "\t".join(cells).rstrip()
            if line.strip():
                yield line
        return
    text = str(getattr(shape, "text", "") or "").strip()
    if text:
        yield text


def extract_pptx_text(
    path: Path,
    *,
    pages: str = "",
    max_slides: int = MAX_DOCUMENT_UNITS_PER_READ,
    max_chars: int = MAX_DOCUMENT_TEXT_CHARS,
) -> PagedTextExtraction:
    """Extract a bounded PPTX slide range, including grouped shapes and tables."""
    source = _inspect_document_file(path)
    _validate_office_archive(source)
    try:
        from pptx import Presentation

        presentation = Presentation(source)
    except Exception as exc:
        raise DocumentReadError(f"Could not open PPTX: {exc}") from exc

    total_slides = len(presentation.slides)
    if total_slides == 0:
        return PagedTextExtraction("", False, 0, 0, -1)
    start, requested_end = _parse_unit_range(pages, total_slides, unit_name="slide")
    end = min(requested_end, start + max(1, int(max_slides)) - 1)
    collector = _TextCollector(max_chars)
    actual_end = start - 1
    for index in range(start, end + 1):
        actual_end = index
        slide_started = False
        for shape in presentation.slides[index].shapes:
            for text in _iter_pptx_shape_text(shape):
                if not slide_started:
                    if not collector.add(f"--- Slide {index + 1} ---", separator="\n\n"):
                        break
                    slide_started = True
                if not collector.add(text, separator="\n"):
                    break
            if collector.truncated:
                break
        if collector.truncated:
            break

    return PagedTextExtraction(
        collector.render(),
        collector.truncated or actual_end < requested_end,
        total_slides,
        start,
        actual_end,
    )


def _resolve_virtual_document(files_root: Path, file_path: str) -> tuple[Path, str, str]:
    raw = str(file_path or "").strip().replace("\\", "/")
    if not raw.startswith("/"):
        raise DocumentReadError(
            "file_path must be an absolute virtual workspace path such as /literature/paper.pdf."
        )
    virtual = PurePosixPath(raw)
    if ".." in virtual.parts:
        raise DocumentReadError("file_path cannot contain parent-directory traversal.")
    relative_parts = list(virtual.parts[1:])
    if relative_parts[:1] == ["files"]:
        relative_parts = relative_parts[1:]
    root = Path(files_root).expanduser().resolve()
    candidate = root.joinpath(*relative_parts).resolve()
    if candidate != root and root not in candidate.parents:
        raise DocumentReadError("file_path resolves outside the workspace files root.")
    suffix = candidate.suffix.lower()
    if suffix not in SUPPORTED_DOCUMENT_SUFFIXES:
        raise DocumentReadError("read_document supports PDF, DOCX, XLSX, and PPTX files only.")
    normalized = "/" + "/".join(relative_parts)
    return candidate, normalized, suffix


def _paged_document_result(
    *,
    label: str,
    virtual_path: str,
    extraction: PagedTextExtraction,
    unit_label: str,
) -> str:
    if extraction.total_units == 0:
        return f"{label} `{virtual_path}` has no {unit_label}s."
    shown_start = extraction.start_unit + 1
    shown_end = extraction.end_unit + 1
    if not extraction.text:
        detail = " For visual inspection, render only the required pages to images." if label == "PDF" else ""
        return (
            f"{label} `{virtual_path}` has no extractable text in {unit_label}s "
            f"{shown_start}-{shown_end}.{detail}"
        )

    result = (
        f"{label} source: `{virtual_path}`\n"
        f"{unit_label.title()}s shown: {shown_start}-{shown_end} of {extraction.total_units}\n\n"
        f"{extraction.text}"
    )
    if extraction.truncated or shown_end < extraction.total_units:
        next_start = min(shown_end + 1, extraction.total_units)
        next_end = min(next_start + MAX_DOCUMENT_UNITS_PER_READ - 1, extraction.total_units)
        if next_start <= extraction.total_units:
            result += f"\n\n(Use pages='{next_start}-{next_end}' to continue.)"
    return result


def read_document(
    files_root: Path,
    *,
    file_path: str,
    pages: str = "",
    cursor: str = "",
) -> str:
    """Return bounded PDF or Office text directly to the caller."""
    try:
        source, virtual_path, suffix = _resolve_virtual_document(files_root, file_path)
        if suffix == ".pdf":
            if str(cursor or "").strip():
                raise DocumentReadError(
                    "cursor applies only to DOCX and XLSX; use pages for PDF."
                )
            return _paged_document_result(
                label="PDF",
                virtual_path=virtual_path,
                extraction=extract_pdf_text(source, pages=pages),
                unit_label="page",
            )
        if suffix == ".pptx":
            if str(cursor or "").strip():
                raise DocumentReadError(
                    "cursor applies only to DOCX and XLSX; use pages for PPTX."
                )
            return _paged_document_result(
                label="PPTX",
                virtual_path=virtual_path,
                extraction=extract_pptx_text(source, pages=pages),
                unit_label="slide",
            )
        if str(pages or "").strip():
            raise DocumentReadError("pages applies only to PDF pages and PPTX slides; leave it empty for DOCX or XLSX.")
        extraction = (
            extract_docx_text(
                source,
                cursor=cursor,
                virtual_path=virtual_path,
            )
            if suffix == ".docx"
            else extract_xlsx_text(
                source,
                cursor=cursor,
                virtual_path=virtual_path,
            )
        )
    except DocumentReadError as exc:
        return f"Error reading document: {exc}"
    except Exception as exc:
        return f"Error reading document: Unexpected parser failure: {exc}"

    label = suffix[1:].upper()
    if not extraction.text:
        return f"{label} `{virtual_path}` has no extractable text."
    result = (
        f"{label} source: `{virtual_path}`\n"
        f"Units shown: {extraction.start_unit_path or 'none'} -> "
        f"{extraction.end_unit_path or 'none'}\n\n{extraction.text}"
    )
    if extraction.next_cursor:
        result += (
            "\n\nThis is a partial document page. Continue with "
            f"cursor='{extraction.next_cursor}'."
        )
    return result


def _tool_call_parts(request: Any) -> tuple[str, str, dict[str, Any]]:
    tool_call = getattr(request, "tool_call", None)
    if not isinstance(tool_call, dict):
        return "", "", {}
    return (
        str(tool_call.get("name") or "").strip(),
        str(tool_call.get("id") or "").strip(),
        dict(tool_call.get("args") or {}) if isinstance(tool_call.get("args"), dict) else {},
    )


def _block_document_identity(block: Any) -> tuple[str, str]:
    if not isinstance(block, dict):
        return "", ""
    source = block.get("source") if isinstance(block.get("source"), dict) else {}
    mime = str(
        block.get("mime_type")
        or block.get("media_type")
        or source.get("mime_type")
        or source.get("media_type")
        or ""
    ).strip().lower()
    filename = str(block.get("filename") or block.get("file_name") or block.get("path") or "").strip()
    return mime, filename


def _is_inline_document_block(block: Any) -> bool:
    if not isinstance(block, dict):
        return False
    mime, filename = _block_document_identity(block)
    suffix = Path(filename).suffix.lower()
    if mime not in DOCUMENT_MIME_TYPES and suffix not in SUPPORTED_DOCUMENT_SUFFIXES:
        return False
    return _has_inline_payload(block)


def _has_inline_payload(block: Any) -> bool:
    if not isinstance(block, dict):
        return False
    source = block.get("source") if isinstance(block.get("source"), dict) else {}
    return any(block.get(key) for key in ("base64", "data")) or any(
        source.get(key) for key in ("base64", "data")
    )


def _document_block_present(message: BaseMessage) -> bool:
    content = message.content
    if not isinstance(content, list):
        return False
    if any(_is_inline_document_block(block) for block in content):
        return True
    if not isinstance(message, ToolMessage) or message.name != "read_file":
        return False
    path = str(message.additional_kwargs.get("read_file_path") or "").strip()
    mime = str(message.additional_kwargs.get("read_file_media_type") or "").strip().lower()
    is_document = Path(path).suffix.lower() in SUPPORTED_DOCUMENT_SUFFIXES or mime in DOCUMENT_MIME_TYPES
    return is_document and any(_has_inline_payload(block) for block in content)


def _document_path_from_message(message: ToolMessage) -> str:
    return str(message.additional_kwargs.get("read_file_path") or "unknown document").strip()


def sanitize_document_tool_message(message: ToolMessage) -> ToolMessage:
    """Replace inline document bytes while preserving message and tool-call identity."""
    if message.name != "read_file" or not _document_block_present(message):
        return message
    path = _document_path_from_message(message)
    mime = str(message.additional_kwargs.get("read_file_media_type") or mimetypes.guess_type(path)[0] or "").strip()
    notice = (
        f"Document binary payload omitted from model history. Source: `{path}` ({mime or 'unknown MIME'}). "
        "Do not retry `read_file` on this document. Use `read_document` for bounded text extraction. "
        "For visual PDF content, render only the required pages to images."
    )
    return message.model_copy(
        update={
            "content": notice,
            "additional_kwargs": {
                "read_file_path": path,
                "read_file_media_type": mime,
                "catmaster_document_payload_removed": True,
            },
        }
    )


def sanitize_document_message(message: BaseMessage) -> BaseMessage:
    """Remove inline PDF or Office bytes from replay before a model call."""
    if isinstance(message, ToolMessage):
        return sanitize_document_tool_message(message)
    if not _document_block_present(message) or not isinstance(message.content, list):
        return message

    content: list[Any] = []
    for block in message.content:
        if not _is_inline_document_block(block):
            content.append(block)
            continue
        _, filename = _block_document_identity(block)
        content.append(
            {
                "type": "text",
                "text": (
                    f"[Inline document payload omitted from replay: {filename or 'stored document'}. "
                    "Use `read_document` with the stored workspace path for bounded text extraction.]"
                ),
            }
        )
    return message.model_copy(update={"content": content})


def _blocked_document_tool_message(*, tool_call_id: str, file_path: str) -> ToolMessage:
    path = str(file_path or "unknown document").strip() or "unknown document"
    mime = mimetypes.guess_type(path)[0] or "application/octet-stream"
    return ToolMessage(
        content=(
            "Direct PDF/Office reads through `read_file` are disabled because raw document bytes would be stored "
            f"in conversation history. Use `read_document(file_path={path!r}, pages='')` for bounded text extraction. "
            "For visual PDF content, render only selected pages to images. Do not retry `read_file` on this document."
        ),
        additional_kwargs={
            "read_file_path": path,
            "read_file_media_type": mime,
            "catmaster_document_read_blocked": True,
        },
        tool_call_id=tool_call_id or "read_file_document_blocked",
        name="read_file",
        status="error",
    )


class DocumentAccessMiddleware(AgentMiddleware):
    """Provide bounded document text access and keep document bytes out of agent state."""

    def __init__(self, *, files_root: Path) -> None:
        self.files_root = Path(files_root).expanduser().resolve()

        def _read_document(
            file_path: str,
            pages: str = "",
            cursor: str = "",
        ) -> str:
            return read_document(
                self.files_root,
                file_path=file_path,
                pages=pages,
                cursor=cursor,
            )

        async def _aread_document(
            file_path: str,
            pages: str = "",
            cursor: str = "",
        ) -> str:
            return await asyncio.to_thread(
                read_document,
                self.files_root,
                file_path=file_path,
                pages=pages,
                cursor=cursor,
            )

        self.tools = [
            StructuredTool.from_function(
                func=_read_document,
                coroutine=_aread_document,
                name="read_document",
                description=(
                    "Read bounded text directly from a workspace PDF, DOCX, XLSX, or PPTX. The result contains "
                    "parsed text and tables, never document bytes. Use pages for PDF pages or PPTX slides. "
                    "For DOCX/XLSX, leave pages empty and pass the returned cursor to continue through unit paths. "
                    "For visual-only PDF content, render selected pages to images."
                ),
                args_schema=ReadDocumentInput,
                infer_schema=False,
            )
        ]

    @property
    def name(self) -> str:
        return "catmaster_document_access"

    @staticmethod
    def _preempt_document_read(request: Any) -> ToolMessage | None:
        tool_name, tool_call_id, args = _tool_call_parts(request)
        file_path = str(args.get("file_path") or "").strip()
        if tool_name == "read_file" and Path(file_path).suffix.lower() in SUPPORTED_DOCUMENT_SUFFIXES:
            return _blocked_document_tool_message(tool_call_id=tool_call_id, file_path=file_path)
        return None

    def wrap_tool_call(self, request: Any, handler: Any) -> Any:
        if blocked := self._preempt_document_read(request):
            return blocked
        result = handler(request)
        if isinstance(result, ToolMessage):
            return sanitize_document_tool_message(result)
        return result

    async def awrap_tool_call(self, request: Any, handler: Any) -> Any:
        if blocked := self._preempt_document_read(request):
            return blocked
        result = await handler(request)
        if isinstance(result, ToolMessage):
            return sanitize_document_tool_message(result)
        return result

    def wrap_model_call(self, request: Any, handler: Any) -> Any:
        messages = [
            sanitize_document_message(message) if isinstance(message, BaseMessage) else message
            for message in request.messages
        ]
        changed = any(updated is not original for updated, original in zip(messages, request.messages))
        sanitized = request.override(messages=messages) if changed else request
        return handler(sanitized)

    async def awrap_model_call(self, request: Any, handler: Any) -> Any:
        messages = [
            sanitize_document_message(message) if isinstance(message, BaseMessage) else message
            for message in request.messages
        ]
        changed = any(updated is not original for updated, original in zip(messages, request.messages))
        sanitized = request.override(messages=messages) if changed else request
        return await handler(sanitized)


__all__ = [
    "BoundedTextExtraction",
    "DOCUMENT_MIME_TYPES",
    "DocumentAccessMiddleware",
    "DocumentReadError",
    "MAX_DOCUMENT_FILE_SIZE_BYTES",
    "MAX_DOCUMENT_TEXT_CHARS",
    "MAX_DOCUMENT_UNITS_PER_READ",
    "PagedTextExtraction",
    "ReadDocumentInput",
    "SUPPORTED_DOCUMENT_SUFFIXES",
    "extract_docx_text",
    "extract_pdf_text",
    "extract_pptx_text",
    "extract_xlsx_text",
    "read_document",
    "sanitize_document_message",
    "sanitize_document_tool_message",
]
